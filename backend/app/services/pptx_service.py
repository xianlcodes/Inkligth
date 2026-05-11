import io
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

TEAL_COLOR = RGBColor(0x0D, 0x94, 0x88)
DARK_COLOR = RGBColor(0x1A, 0x1A, 0x2E)
WHITE_COLOR = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF0, 0xFD, 0xFB)


def generate_pptx(outline: dict, paper_title: str = "") -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides = outline.get("slides", [])

    for i, slide_data in enumerate(slides):
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout

        # Background
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = WHITE_COLOR

        # Top accent bar
        bar = slide.shapes.add_shape(
            1, Inches(0), Inches(0), prs.slide_width, Inches(0.06)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = TEAL_COLOR
        bar.line.fill.background()

        # Slide number
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(2), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = f"{i + 1:02d}"
        p.font.size = Pt(14)
        p.font.color.rgb = TEAL_COLOR
        p.font.bold = True

        # Title
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = DARK_COLOR

        # Title underline
        line = slide.shapes.add_shape(
            1, Inches(0.8), Inches(1.55), Inches(3.5), Inches(0.04)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = TEAL_COLOR
        line.line.fill.background()

        # Bullets
        bullets = slide_data.get("bullets", [])
        if bullets:
            bullet_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(4.5))
            tf = bullet_box.text_frame
            tf.word_wrap = True

            for j, bullet in enumerate(bullets):
                if j == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = f"• {bullet}"
                p.font.size = Pt(20)
                p.font.color.rgb = DARK_COLOR
                p.space_after = Pt(14)
                p.line_spacing = Pt(30)

        # Notes
        notes_text = slide_data.get("notes", "")
        if notes_text:
            notes_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.6))
            tf = notes_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = f"💡 {notes_text}"
            p.font.size = Pt(12)
            p.font.color.rgb = TEAL_COLOR
            p.font.italic = True

        # Footer line
        footer_line = slide.shapes.add_shape(
            1, Inches(0), Inches(7.3), prs.slide_width, Inches(0.02)
        )
        footer_line.fill.solid()
        footer_line.fill.fore_color.rgb = TEAL_COLOR
        footer_line.line.fill.background()

        # Footer text
        footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.32), Inches(12.3), Inches(0.18))
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = paper_title or "Literature Report"
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
        p.alignment = PP_ALIGN.RIGHT

    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    logger.info(f"PPTX generated with {len(slides)} slides")
    return buffer.getvalue()