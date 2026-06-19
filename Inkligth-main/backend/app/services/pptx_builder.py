"""专业 PPT 渲染引擎 — 替代旧的 pptx_service.py

支持学科主题、多布局模板、图像嵌入、公式渲染、图表生成。
"""

import io
import logging
import os
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

from app.services.ppt_theme import DisciplineTheme, get_theme

logger = logging.getLogger(__name__)

# ── 常量 ──
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

MARGIN_LR = Inches(0.8)
MARGIN_TB = Inches(0.5)

CONTENT_WIDTH = SLIDE_WIDTH - MARGIN_LR * 2
CONTENT_HEIGHT = Inches(4.5)

TITLE_TOP = Inches(0.6)
TITLE_LEFT = MARGIN_LR
TITLE_WIDTH = CONTENT_WIDTH
TITLE_HEIGHT = Inches(1.0)

BODY_TOP = Inches(2.0)
BODY_LEFT = MARGIN_LR

ACCENT_BAR_H = Inches(0.06)
FOOTER_TOP = Inches(7.3)
FOOTER_H = Inches(0.2)


def _hex_to_rgb(hex_color: str) -> RGBColor:
    """#RRGGBB → RGBColor"""
    hex_color = hex_color.lstrip("#")
    return RGBColor(
        int(hex_color[0:2], 16),
        int(hex_color[2:4], 16),
        int(hex_color[4:6], 16),
    )


def _add_textbox(slide, left: float, top: float, width: float, height: float):
    """添加强制指定尺寸的文本框（Emu 转换）"""
    return slide.shapes.add_textbox(
        int(left), int(top), int(width), int(height)
    )


def _set_font(run, size_pt: int = 18, bold: bool = False,
              color: str = "#1A1A2E", italic: bool = False,
              font_name: str = ""):
    """设置 run 的字体属性"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = _hex_to_rgb(color)
    if font_name:
        run.font.name = font_name
        # 东亚字体
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn("a:ea"), font_name)


def _set_shape_fill(shape, color: str, alpha: Optional[int] = None):
    """设置形状填充色"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = _hex_to_rgb(color)
    shape.line.fill.background()


class PptxBuilder:
    """专业 PPT 渲染引擎"""

    def __init__(self, theme: DisciplineTheme = None):
        self.theme = theme or get_theme("cs")

    def build(self, slides_data: list[dict],
              visual_assets: Optional[dict] = None,
              paper_title: str = "") -> bytes:
        """
        根据增强大纲和视觉资产渲染 PPTX。

        slides_data: OutlineGenerator 输出的 slide 列表（含 page_type/visual_ref）
        visual_assets: VisualExtractor 输出的资产字典
        paper_title: 论文标题（用于页脚）
        """
        prs = Presentation()
        prs.slide_width = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        for i, slide_data in enumerate(slides_data):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

            page_type = slide_data.get("page_type", "bullet_text")

            # ── 根据 page_type 渲染 ──
            if page_type == "title":
                self._render_title_slide(slide, slide_data)
            elif page_type == "section_header":
                self._render_section_header(slide, slide_data)
            elif page_type in ("figure_full", "figure_text"):
                self._render_with_image(slide, slide_data, page_type,
                                        visual_assets)
            elif page_type == "formula_derivation":
                self._render_formula_slide(slide, slide_data, visual_assets)
            elif page_type == "data_chart":
                self._render_chart_slide(slide, slide_data, paper_title)
            else:
                self._render_default_slide(slide, slide_data)

            # ── 共用装饰 ──
            self._add_decorations(slide, i + 1, len(slides_data),
                                  slide_data, paper_title)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        logger.info(f"PPTX generated with {len(slides_data)} slides "
                    f"(theme={self.theme.name})")
        return buffer.getvalue()

    # ── 各 page_type 渲染方法 ──

    def _render_title_slide(self, slide, data: dict):
        """封面页"""
        # 渐变背景色块
        bg_shape = slide.shapes.add_shape(
            1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
        )
        _set_shape_fill(bg_shape, self.theme.primary_color)
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = _hex_to_rgb(self.theme.primary_color)

        # 标题
        title = data.get("title", "")
        tb = _add_textbox(slide, Inches(1.5), Inches(2.0),
                          Inches(10.333), Inches(2.5))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        _set_font(p.runs[0] if p.runs else p,
                  size_pt=40, bold=True,
                  color=self.theme.text_light)

        # 副标题/作者信息
        bullets = data.get("bullets", [])
        if bullets:
            sub_tb = _add_textbox(slide, Inches(2), Inches(4.5),
                                  Inches(9.333), Inches(1.5))
            sub_tb.text_frame.word_wrap = True
            for j, b in enumerate(bullets):
                p = sub_tb.text_frame.paragraphs[0] if j == 0 else \
                    sub_tb.text_frame.add_paragraph()
                p.text = b
                p.alignment = PP_ALIGN.CENTER
                _set_font(p.runs[0] if p.runs else p,
                          size_pt=18, color=self.theme.text_light)

    def _render_section_header(self, slide, data: dict):
        """章节过渡页"""
        bg_shape = slide.shapes.add_shape(
            1, 0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
        )
        _set_shape_fill(bg_shape, self.theme.secondary_color)

        title = data.get("title", "")
        tb = _add_textbox(slide, Inches(1.5), Inches(2.5),
                          Inches(10.333), Inches(2.0))
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.alignment = PP_ALIGN.CENTER
        _set_font(p.runs[0] if p.runs else p,
                  size_pt=38, bold=True, color=self.theme.text_light)

    def _render_default_slide(self, slide, data: dict):
        """默认布局：标题 + 要点"""
        self._add_title(slide, data.get("title", ""))
        bullets = data.get("bullets", [])
        if bullets:
            self._add_bullets(slide, bullets)

    def _render_with_image(self, slide, data: dict, page_type: str,
                           visual_assets: Optional[dict]):
        """含图片的布局：figure_full / figure_text"""
        self._add_title(slide, data.get("title", ""))

        ref = data.get("visual_ref")
        image_path = None
        if ref and visual_assets:
            for fig in visual_assets.get("figures", []):
                if fig["id"] == ref and fig.get("image_path") \
                        and os.path.exists(fig["image_path"]):
                    image_path = fig["image_path"]
                    break

        bullets = data.get("bullets", [])

        if page_type == "figure_full":
            # 全页图（在标题下方）
            if image_path:
                slide.shapes.add_picture(
                    image_path, BODY_LEFT, BODY_TOP,
                    CONTENT_WIDTH, CONTENT_HEIGHT
                )
        else:
            # 左图右文
            if image_path:
                img_w = Inches(6.0)
                slide.shapes.add_picture(
                    image_path, BODY_LEFT, BODY_TOP,
                    img_w, CONTENT_HEIGHT
                )
            if bullets:
                self._add_bullets(
                    slide, bullets,
                    left=BODY_LEFT + Inches(6.5),
                    width=CONTENT_WIDTH - Inches(6.5),
                )

    def _render_formula_slide(self, slide, data: dict,
                               visual_assets: Optional[dict]):
        """公式页：公式在上，解释在下"""
        self._add_title(slide, data.get("title", ""))

        ref = data.get("visual_ref")
        formula_image = None
        if ref and visual_assets:
            for f in visual_assets.get("formulas", []):
                if f["id"] == ref and f.get("image_path") \
                        and os.path.exists(f["image_path"]):
                    formula_image = f["image_path"]
                    break

        if formula_image:
            slide.shapes.add_picture(
                formula_image, BODY_LEFT, BODY_TOP,
                CONTENT_WIDTH, Inches(2.5)
            )

        bullets = data.get("bullets", [])
        if bullets:
            self._add_bullets(
                slide, bullets,
                top=BODY_TOP + Inches(3.0),
                height=Inches(3.0),
            )

    def _render_chart_slide(self, slide, data: dict, paper_title: str):
        """图表页：自动生成的图表 + 说明文字"""
        from app.services.chart_renderer import ChartRenderer

        self._add_title(slide, data.get("title", ""))

        chart_type = data.get("suggested_chart", "bar")
        data_hint = data.get("chart_data_hint", "")

        if chart_type and data_hint:
            cache_dir = "/tmp/charts"
            renderer = ChartRenderer(cache_dir, {
                "primary": self.theme.primary_color,
                "secondary": self.theme.secondary_color,
                "accent": self.theme.accent_color,
            })
            chart_path = renderer.render(
                chart_type, data_hint, self.theme.name
            )
            if chart_path and os.path.exists(chart_path):
                slide.shapes.add_picture(
                    chart_path, BODY_LEFT, BODY_TOP,
                    CONTENT_WIDTH, Inches(3.5)
                )

        bullets = data.get("bullets", [])
        if bullets:
            self._add_bullets(
                slide, bullets,
                top=BODY_TOP + Inches(3.8),
                height=Inches(2.0),
                font_size=14,
            )

    # ── 辅助渲染方法 ──

    def _add_title(self, slide, title: str):
        """添加幻灯片标题"""
        if not title:
            return
        tb = _add_textbox(slide, TITLE_LEFT, TITLE_TOP,
                          TITLE_WIDTH, TITLE_HEIGHT)
        tb.text_frame.word_wrap = True
        p = tb.text_frame.paragraphs[0]
        p.text = title
        _set_font(p.runs[0] if p.runs else p,
                  size_pt=self.theme.font_sizes["title"],
                  bold=True, color=self.theme.text_primary)

        # 标题下装饰线
        line = slide.shapes.add_shape(
            1, TITLE_LEFT, Inches(1.55),
            Inches(3.5), Inches(0.04)
        )
        _set_shape_fill(line, self.theme.primary_color)

    def _add_bullets(self, slide, bullets: list,
                     left=None, top=None, width=None, height=None,
                     font_size=None):
        """添加要点正文"""
        left = left or BODY_LEFT
        top = top or BODY_TOP
        width = width or CONTENT_WIDTH
        height = height or CONTENT_HEIGHT
        font_size = font_size or self.theme.font_sizes["body"]

        tb = _add_textbox(slide, left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        for j, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.text = f"• {bullet}"
            p.space_after = Pt(8)
            p.line_spacing = Pt(font_size * 1.5)
            _set_font(p.runs[0] if p.runs else p,
                      size_pt=font_size, color=self.theme.text_primary)

    def _add_decorations(self, slide, slide_num: int, total: int,
                         data: dict, paper_title: str):
        """添加共用装饰元素（编号、页脚、演讲者备注）"""
        # 顶部装饰条
        if self.theme.accent_bar:
            bar = slide.shapes.add_shape(
                1, 0, 0, SLIDE_WIDTH, ACCENT_BAR_H
            )
            _set_shape_fill(bar, self.theme.primary_color)

        # 幻灯片编号
        num_box = _add_textbox(slide, Inches(0.5), Inches(0.2),
                               Inches(2), Inches(0.4))
        p = num_box.text_frame.paragraphs[0]
        p.text = f"{slide_num:02d}"
        _set_font(p, size_pt=14, bold=True, color=self.theme.secondary_color)

        # 底部装饰线
        footer_line = slide.shapes.add_shape(
            1, 0, FOOTER_TOP, SLIDE_WIDTH, Inches(0.02)
        )
        _set_shape_fill(footer_line, self.theme.divider_color)

        # 底部文字
        footer_box = _add_textbox(slide, Inches(0.5), Inches(7.32),
                                  Inches(12.3), FOOTER_H)
        p = footer_box.text_frame.paragraphs[0]
        p.text = paper_title or "Literature Report"
        p.alignment = PP_ALIGN.RIGHT
        _set_font(p, size_pt=9, color=self.theme.divider_color)

        # 演讲者备注
        notes_text = data.get("notes", "") or data.get("speaker_notes", "")
        if notes_text:
            notes_slide = slide.notes_slide
            tf = notes_slide.notes_text_frame
            p = tf.paragraphs[0]
            p.text = notes_text
